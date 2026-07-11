import io

from pptx import Presentation
from pptx.util import Inches


def generar_presentacion_curso(data) -> bytes:
    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    datos = data.datos
    objetivos = data.objetivos
    temario = data.temario
    encuadre = data.encuadre
    evaluaciones = data.evaluaciones

    beneficios = data.beneficios or ""
    expositiva = data.expositiva or {}
    demostrativa = data.demostrativa or {}
    dialogo = data.dialogo or {}
    cierre = data.cierre or {}

    def texto_seguro(valor):
        if valor is None:
            return ""
        if isinstance(valor, list):
            return "\n".join(str(v) for v in valor if v)
        return str(valor)

    def agregar_slide(titulo, contenido="", notas=""):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = texto_seguro(titulo)
        body = slide.placeholders[1]
        body.text = texto_seguro(contenido)
        if notas:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = texto_seguro(notas)
        return slide

    def agregar_portada():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = datos.nombreCurso or "Nombre del taller"
        subtitle = slide.placeholders[1]
        subtitle.text = (
            f"Instructor: {datos.instructor or ''}\n"
            f"Fecha: {datos.fecha or ''}\n"
            f"Lugar: {datos.lugar or ''}"
        )

    def lista_temario():
        lineas = []
        if getattr(temario, "u1", None):
            lineas.append("Unidad 1:")
            lineas.extend([f"• {t}" for t in temario.u1])
        if getattr(temario, "u2", None):
            lineas.append("\nUnidad 2:")
            lineas.extend([f"• {t}" for t in temario.u2])
        if getattr(temario, "u3", None):
            lineas.append("\nUnidad 3:")
            lineas.extend([f"• {t}" for t in temario.u3])
        return "\n".join(lineas)

    def texto_encuadre_reglas():
        reglas = []
        if hasattr(encuadre, "reglasTexto") and encuadre.reglasTexto:
            reglas.extend(encuadre.reglasTexto)
        if hasattr(encuadre, "otraRegla") and encuadre.otraRegla:
            reglas.append(encuadre.otraRegla)
        return "\n".join([f"• {r}" for r in reglas])

    def texto_encuadre_contrato():
        acuerdos = []
        if hasattr(encuadre, "acuerdosTexto") and encuadre.acuerdosTexto:
            acuerdos.extend(encuadre.acuerdosTexto)
        if hasattr(encuadre, "otroAcuerdo") and encuadre.otroAcuerdo:
            acuerdos.append(encuadre.otroAcuerdo)
        return "\n".join([f"• {a}" for a in acuerdos])

    def texto_preguntas_experiencia():
        return encuadre.preguntas or "" if hasattr(encuadre, "preguntas") else ""

    agregar_portada()
    agregar_slide("Presentación del Ponente", (
        f"Nombre del ponente: {datos.instructor or ''}\n\n"
        "Breve descripción de su formación profesional:\n\n"
        "Áreas de interés:\n\n"
        "Datos de contacto:\n\n"
        "Horario y lugar donde lo pueden localizar:\n\n"
        "Fotografía:"
    ))
    agregar_slide("Lo que veremos en este curso",
        f"Curso: {datos.nombreCurso or ''}\n\nPerfil de participantes:\n{datos.perfil or ''}")
    agregar_slide("Objetivo General", objetivos.general or "")
    agregar_slide("Objetivos Particulares", (
        f"Cognitivo:\n{objetivos.cognitiva or ''}\n\n"
        f"Psicomotriz:\n{objetivos.psicomotriz or ''}\n\n"
        f"Afectivo:\n{objetivos.afectiva or ''}"
    ))
    agregar_slide("Temario", lista_temario())
    agregar_slide("Beneficios de este curso", beneficios)
    agregar_slide("¿Cómo será la evaluación del curso?", (
        "Evaluación diagnóstica\n"
        f"Evaluación formativa: {evaluaciones.pctFormativa or 0}%\n"
        f"Evaluación sumativa: {evaluaciones.pctSumativa or 0}%\n"
        "Evaluación de reacción"
    ))
    agregar_slide("¿Qué esperas de este curso?", texto_preguntas_experiencia())
    agregar_slide("Reglas del curso", texto_encuadre_reglas())
    agregar_slide("Contrato de aprendizaje", texto_encuadre_contrato())
    agregar_slide("Evaluación Diagnóstica", f"Instrucciones:\n{evaluaciones.instDiagnostica or ''}")
    agregar_slide("DESARROLLO DE LOS TEMAS",
        "A partir de aquí se deberán incluir todas las diapositivas de la técnica expositiva.")
    agregar_slide("Fin de la técnica expositiva", "")
    agregar_slide("Actividad", (
        f"Nombre de la actividad:\n{demostrativa.get('actividad', '')}\n\n"
        f"Instrucciones:\n{demostrativa.get('ejemplos', '')}"
    ))
    agregar_slide("Evaluación Formativa", (
        f"Instrucciones:\n{evaluaciones.instFormativa or ''}\n\n"
        f"Ponderación: {evaluaciones.pctFormativa or 0}%"
    ))
    agregar_slide("Debate", "Instrucciones:")
    agregar_slide("Diálogo, foro, mesa redonda, etc.",
        f"Instrucciones:\n{dialogo.get('instrucciones', '')}")
    agregar_slide("Evaluación Final", (
        f"Instrucciones:\n{evaluaciones.instSumativa or ''}\n\n"
        f"Ponderación: {evaluaciones.pctSumativa or 0}%"
    ))
    agregar_slide("Resumen", cierre.get('resumen', ''))
    agregar_slide("Logro de expectativas", "")
    agregar_slide("Logro de objetivos", "")
    agregar_slide("Sugerencias de continuidad del aprendizaje", cierre.get('sugerencias', ''))
    agregar_slide("Compromisos de aplicación del aprendizaje", cierre.get('compromisos', ''))
    agregar_slide("Ayúdame a llenar la siguiente encuesta",
        f"Instrucciones:\n{evaluaciones.instReac or ''}")
    agregar_slide("¡GRACIAS!", "")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
