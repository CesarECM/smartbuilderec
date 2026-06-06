from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from starlette.middleware.base import BaseHTTPMiddleware
from pydantic import BaseModel
from typing import Optional
from dotenv import load_dotenv
from openai import OpenAI
from pathlib import Path
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches
import os
import json
import asyncio
import pyotp
import zipfile
import io
import subprocess
import base64
from concurrent.futures import ThreadPoolExecutor, as_completed
import sys
sys.path.insert(0, str(Path(__file__).resolve().parent))

BASE_DIR = Path(__file__).resolve().parent

load_dotenv()

app = FastAPI()

import base64 as _b64
import json as _json

_PATHS_PUBLICOS = {"/", "/validate-token", "/webhook/stripe", "/checkout/session", "/checkout/verify"}
_jwks_cache: dict = {}

# Endpoints de soporte accesibles sin JWT (widget anónimo desde landing, pago, etc.)
_SOPORTE_PUBLICOS = {
    ("POST", "/soporte/sesiones/init"),
    ("POST", "/soporte/chat"),
    ("POST", "/soporte/tickets"),
    ("POST", "/soporte/cron/analizar-patrones"),
    ("POST", "/admin/cron/vigencias"),      # protegido por CRON_SECRET
    ("POST", "/auth/reset-password"),       # protegido a nivel de lógica (no JWT)
}

def _es_publico(method: str, path: str) -> bool:
    if path in _PATHS_PUBLICOS:
        return True
    if (method, path) in _SOPORTE_PUBLICOS:
        return True
    return False

async def _get_public_key(supabase_url: str, kid: str):
    """Obtiene y cachea la clave pública del JWKS de Supabase."""
    global _jwks_cache
    if kid in _jwks_cache:
        return _jwks_cache[kid]
    try:
        import httpx
        async with httpx.AsyncClient(timeout=5.0) as client:
            res = await client.get(f"{supabase_url}/auth/v1/.well-known/jwks.json")
            if res.status_code == 200:
                for k in res.json().get("keys", []):
                    _jwks_cache[k["kid"]] = k
    except Exception:
        pass
    return _jwks_cache.get(kid)

class JWTAuthMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        if request.method == "OPTIONS" or _es_publico(request.method, request.url.path):
            return await call_next(request)
        auth = request.headers.get("Authorization", "")
        if not auth.startswith("Bearer "):
            return JSONResponse({"detail": "Autenticación requerida."}, status_code=401)
        token = auth.removeprefix("Bearer ").strip()
        try:
            from jose import jwt as jose_jwt, jwk as jose_jwk

            # Leer alg y kid del header del token
            raw_hdr = token.split('.')[0]
            raw_hdr += '=' * (-len(raw_hdr) % 4)
            header = _json.loads(_b64.urlsafe_b64decode(raw_hdr))
            alg = header.get('alg', 'HS256')
            kid = header.get('kid', '')

            if alg.startswith('ES') or alg.startswith('RS') or alg.startswith('PS'):
                # Asimétrico (ES256, RS256…): verificar con clave pública desde JWKS
                supabase_url = os.getenv("SUPABASE_URL", "")
                if not supabase_url:
                    return JSONResponse({"detail": "SUPABASE_URL no configurado."}, status_code=500)
                key_data = await _get_public_key(supabase_url, kid)
                if not key_data:
                    return JSONResponse({"detail": f"Clave pública no encontrada (kid={kid})."}, status_code=401)
                public_key = jose_jwk.construct(key_data, algorithm=alg)
                payload = jose_jwt.decode(token, public_key, algorithms=[alg], options={"verify_aud": False})
            else:
                # Simétrico (HS256…): verificar con SUPABASE_JWT_SECRET
                secret = os.getenv("SUPABASE_JWT_SECRET")
                if not secret:
                    return JSONResponse({"detail": "JWT secret no configurado."}, status_code=500)
                payload = jose_jwt.decode(token, secret, algorithms=[alg], options={"verify_aud": False})

            request.state.user = payload
        except Exception as e:
            return JSONResponse({"detail": f"Token inválido o expirado: {str(e)}"}, status_code=401)
        return await call_next(request)

# IMPORTANTE: JWTAuthMiddleware primero (inner), CORSMiddleware después (outer).
# En Starlette los middlewares se aplican LIFO: el último en add_middleware
# es el más externo y procesa requests primero / responses último.
# CORSMiddleware debe ser outer para añadir headers CORS incluso en respuestas 401/500.
app.add_middleware(JWTAuthMiddleware)
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5500",
        "http://127.0.0.1:5500",
        "https://ceecm-web.vercel.app",
        "https://smartbuilderec.vercel.app",
        "https://smartbuilderec.com",
        "https://www.smartbuilderec.com",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["x-sesion-id", "x-faq-ids"],
)

# ─── Routers ──────────────────────────────────────────────────────────────────

from routers import stripe_router
from routers import admin_router
from routers import email_router
from routers import soporte_router

app.include_router(stripe_router.router,  tags=["stripe"])
app.include_router(admin_router.router,   tags=["admin"])
app.include_router(email_router.router,   tags=["email"])
app.include_router(soporte_router.router, tags=["soporte"])

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o")
OPENAI_MODEL_GRL = os.getenv("OPENAI_MODEL_GRL", "gpt-4o-mini")


# ─── Modelos de datos ─────────────────────────────────────────────────────────

class EvaluationRequest(BaseModel):
    texto: str
    tipo: str
    objetivo_cognitivo: str = ""
    objetivo_psicomotriz: str = ""

class GeneralRequest(BaseModel):
    cognitiva: str
    psicomotriz: str
    afectiva: str

class BeneficiosRequest(BaseModel):
    general: str
    nombre: str = ""

class TemarioRequest(BaseModel):
    nombre: str = ""
    general: str = ""
    cognitiva: str = ""
    psicomotriz: str = ""
    afectiva: str = ""
    beneficios: str = ""

class TokenRequest(BaseModel):
    token: str

class ObjetivosRequest(BaseModel):
    general: str
    cognitiva: str
    psicomotriz: str
    afectiva: str

class DatosInfo(BaseModel):
    nombreCurso: str = ""
    instructor: str = ""
    disenador: str = ""
    lugar: str = ""
    fecha: str = ""
    duracion: int | None = None
    participantes: int | None = None
    perfil: str = ""

class ObjetivosInfo(BaseModel):
    general: str = ""
    cognitiva: str = ""
    psicomotriz: str = ""
    afectiva: str = ""

class TemarioInfo(BaseModel):
    u1: list[str] = []
    u2: list[str] = []
    u3: list[str] = []


class PreguntasRequest(BaseModel):
    nombre: str = ""
    perfil: str = ""

class EncuadreInfo(BaseModel):
    preguntas: str = ""
    reglas: list[int] = []
    reglasTexto: list[str] = []
    otraRegla: str = ""
    acuerdos: list[int] = []
    acuerdosTexto: list[str] = []
    otroAcuerdo: str = ""

class TecnicasInfo(BaseModel):
    rhIdx: int = 0
    rhNombre: str = ""
    rhObjetivo: str = ""
    rhInstrucciones: str = ""
    rhDetalle: str = ""
    rhDuracion: str = ""
    rhMateriales: str = ""
    rhCustom: str = ""
    rhSeleccion: str = ""
    enIdx: int = 0
    enNombre: str = ""
    enObjetivo: str = ""
    enInstrucciones: str = ""
    enDetalle: str = ""
    enDuracion: str = ""
    enMateriales: str = ""
    enCustom: str = ""
    enSeleccion: str = ""
    rompehielos: dict = {}
    energizante: dict = {}

class EvaluacionesInfo(BaseModel):
    pctDiag: int = 0
    pctDiagnostica: int = 0
    pctForm: int = 0
    pctFormativa: int = 0
    pctSuma: int = 0
    pctSumativa: int = 0
    instDiag: str = ""
    instDiagnostica: str = ""
    instForm: str = ""
    instFormativa: str = ""
    instSuma: str = ""
    instSumativa: str = ""
    instReac: str = ""
    descripcionGeneral: str = ""
    tipoInstrumentoFormativa: str = ""

class TiempoFila(BaseModel):
    titulo: str = ""
    tiempo: int = 0

class TiempoBloque(BaseModel):
    seccion: str = ""
    filas: list[TiempoFila] = []

class PlaneacionRequest(BaseModel):
    datos: DatosInfo = DatosInfo()
    objetivos: ObjetivosInfo = ObjetivosInfo()
    beneficios: str = ""
    temario: TemarioInfo = TemarioInfo()
    encuadre: EncuadreInfo = EncuadreInfo()
    tecnicas: TecnicasInfo = TecnicasInfo()
    evaluaciones: EvaluacionesInfo = EvaluacionesInfo()
    tiempos: list[TiempoBloque] = []
    expositiva: dict = {}
    demostrativa: dict = {}
    dialogo: dict = {}
    cierre: dict = {}
    materiales: dict = {}


class ExpositivaRequest(BaseModel):
    campo: str
    nombreCurso: str = ""
    perfil: str = ""
    objetivoCognitivo: str = ""
    objetivoGeneral: str = ""
    temario: dict = {}


class DemostrativaRequest(BaseModel):
    campo: str
    nombreCurso: str = ""
    perfil: str = ""
    objetivoPsicomotriz: str = ""
    objetivoGeneral: str = ""
    temario: dict = {}



class DialogoRequest(BaseModel):
    campo: str
    nombreCurso: str = ""
    perfil: str = ""
    objetivoAfectivo: str = ""
    objetivoGeneral: str = ""
    temario: dict = {}

class CierreRequest(BaseModel):
    nombreCurso: str = ""
    objetivoGeneral: str = ""
    objetivoCognitivo: str = ""
    objetivoPsicomotriz: str = ""
    objetivoAfectivo: str = ""
    desarrolloExpositiva: str = ""
    actividadDemostrativa: str = ""
    instruccionesDialogo: str = ""



class DescripcionGeneralRequest(BaseModel):
    cierre: str = ""


class EvaluacionIARequest(BaseModel):
    nombreCurso: str = ""
    objetivoGeneral: str = ""
    objetivoCognitivo: str = ""
    objetivoPsicomotriz: str = ""
    objetivoAfectivo: str = ""

class MaterialesRequest(BaseModel):
    tecnica: str = ""
    nombreCurso: str = ""
    perfil: str = ""
    objetivoGeneral: str = ""
    objetivos: dict = {}
    temario: dict = {}
    tecnicas: dict = {}
    expositiva: dict = {}
    demostrativa: dict = {}
    dialogo: dict = {}

# ─── Utilidades ───────────────────────────────────────────────────────────────

def load_prompt(filename: str, **kwargs) -> str:
    file_path = BASE_DIR / "prompts" / filename
    if not file_path.exists():
        raise FileNotFoundError(f"No se encontró el archivo {filename}")
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()
    return content.format(**kwargs)


# ─── Generación del .docx de objetivos ───────────────────────────────────────

def crear_docx_objetivos(data: ObjetivosRequest) -> bytes:
    doc = Document()
    titulo = doc.add_heading("", level=1)
    titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_titulo = titulo.add_run("Objetivos de Aprendizaje")
    run_titulo.font.color.rgb = RGBColor(0x1F, 0x3B, 0x6D)
    run_titulo.font.size = Pt(18)
    doc.add_paragraph()
    secciones = [
        ("Objetivo General",     data.general),
        ("Objetivo Cognitivo",   data.cognitiva),
        ("Objetivo Psicomotriz", data.psicomotriz),
        ("Objetivo Afectivo",    data.afectiva),
    ]
    for titulo_sec, texto in secciones:
        h = doc.add_heading("", level=2)
        run_h = h.add_run(titulo_sec)
        run_h.font.color.rgb = RGBColor(0x31, 0x4E, 0x7A)
        run_h.font.size = Pt(14)
        p = doc.add_paragraph()
        run_p = p.add_run(texto)
        run_p.font.size = Pt(12)
        doc.add_paragraph()
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()




def crear_zip_con_docx(docx_bytes: bytes, nombre: str = "objetivos_EC0217.docx") -> bytes:
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(nombre, docx_bytes)
    zip_buffer.seek(0)
    return zip_buffer.read()


# ─── Generación del Documento de Planeación (Node.js) ────────────────────────

def generar_planeacion_docx(payload: dict) -> bytes:
    script_path = BASE_DIR / "generar_planeacion.js"
    if not script_path.exists():
        raise FileNotFoundError("No se encontró generar_planeacion.js en el backend")

    # Normalizar evaluaciones: el frontend puede mandar pctDiag o pctDiagnostica, etc.
    ev = payload.get("evaluaciones", {})
    payload["evaluaciones"] = {
        "pctDiagnostica":          ev.get("pctDiagnostica", ev.get("pctDiag", 0)),
        "instDiagnostica":         ev.get("instDiagnostica", ev.get("instDiag", "")),
        "pctFormativa":            ev.get("pctFormativa",    ev.get("pctForm", 0)),
        "instFormativa":           ev.get("instFormativa",   ev.get("instForm", "")),
        "pctSumativa":             ev.get("pctSumativa",     ev.get("pctSuma", 0)),
        "instSumativa":            ev.get("instSumativa",    ev.get("instSuma", "")),
        "instReac":                ev.get("instReac", ""),
        "descripcionGeneral":      ev.get("descripcionGeneral", ""),
        "tipoInstrumentoFormativa": ev.get("tipoInstrumentoFormativa", ""),
    }

    # Normalizar tecnicas: asegurar que rhDetalle y enDetalle estén presentes.
    # El frontend los construye pero si vienen vacíos los reconstruimos desde objetivo+instrucciones.
    tc = payload.get("tecnicas", {})
    rh_objetivo     = tc.get("rhObjetivo", "")
    rh_instrucciones = tc.get("rhInstrucciones", "")
    en_objetivo     = tc.get("enObjetivo", "")
    en_instrucciones = tc.get("enInstrucciones", "")

    if not tc.get("rhDetalle") and (rh_objetivo or rh_instrucciones):
        tc["rhDetalle"] = (
            f"a) Explicará objetivo de la técnica:\n{rh_objetivo}\n\n"
            f"b) Dará las instrucciones de la técnica:\n{rh_instrucciones}\n\n"
            f"c) Mencionará el tiempo para realizarla.\n\n"
            f"d) Propiciará la participación del grupo.\n\n"
            f"e) Integrará al grupo.\n\n"
            f"f) Controlará el tiempo."
        )
    if not tc.get("enDetalle") and (en_objetivo or en_instrucciones):
        tc["enDetalle"] = (
            f"a) Explicará objetivo de la técnica:\n{en_objetivo}\n\n"
            f"b) Dará las instrucciones de la técnica:\n{en_instrucciones}"
        )
    # Asegurar que enObjetivo también esté en el nivel raíz de tecnicas para el JS
    if not tc.get("enObjetivo") and en_objetivo:
        tc["enObjetivo"] = en_objetivo
    payload["tecnicas"] = tc

    payload_json = json.dumps(payload, ensure_ascii=False)
    result = subprocess.run(
        ["node", str(script_path)],
        input=payload_json.encode("utf-8"),
        capture_output=True,
        cwd=str(BASE_DIR),
        timeout=30
    )

    stderr_msg = result.stderr.decode("utf-8", errors="replace").strip()
    if result.returncode != 0 or not result.stdout:
        raise RuntimeError(f"Error al generar el documento: {stderr_msg or 'stdout vacío'}")

    try:
        return base64.b64decode(result.stdout)
    except Exception as e:
        raise RuntimeError(f"Error al decodificar base64: {e}. stderr: {stderr_msg}")
    



def generar_evaluacion_diagnostica(data: PlaneacionRequest) -> bytes:
    doc = Document()

    titulo = doc.add_heading("EVALUACIÓN DIAGNÓSTICA", level=1)
    titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitulo = doc.add_heading("CUESTIONARIO", level=2)
    subtitulo.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()

    tabla_encabezado(doc, data.datos)

    # ── Instrucciones de aplicación ───────────────────────────────────────────
    doc.add_heading("Instrucciones de aplicación", level=2)
    instrucciones = [
        ("Alcance/beneficio:", "El presente cuestionario tiene como propósito conocer los conocimientos previos del participante antes de iniciar el curso."),
        ("Indicaciones para el facilitador:", "Entregue el cuestionario al participante antes de iniciar el curso. No hay respuestas correctas o incorrectas; el objetivo es identificar el nivel de conocimiento previo."),
        ("Condiciones de aplicación:", "Se aplicará al inicio del curso, de manera individual y sin consultas documentales."),
        ("Tiempo para desarrollar la actividad:", "Considerar el tiempo establecido en el documento de planeación."),
        ("Valor:", f"{data.evaluaciones.pctDiagnostica or data.evaluaciones.pctDiag or 0}% de la calificación total."),
    ]
    for etiqueta, texto in instrucciones:
        p = doc.add_paragraph()
        p.add_run(etiqueta).bold = True
        p.add_run(f" {texto}")

    doc.add_paragraph()

    # ── Reactivos generados por IA ────────────────────────────────────────────
    reactivos_texto = (
        data.evaluaciones.instDiagnostica
        or data.evaluaciones.instDiag
        or ""
    ).strip()

    if reactivos_texto:
        doc.add_heading("Cuestionario", level=2)
        for linea in reactivos_texto.split("\n"):
            if linea.strip():
                doc.add_paragraph(linea.strip())
    else:
        doc.add_heading("Cuestionario", level=2)
        for i in range(1, 11):
            doc.add_paragraph(f"{i}. __________________________________________________________")
            doc.add_paragraph()

    doc.add_paragraph()
    doc.add_paragraph("Nombre y firma del Participante: ________________________________________")

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    return buffer.read()


def generar_presentacion_curso(data) -> bytes:
    prs = Presentation()

    # Tamaño widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    datos = data.datos
    objetivos = data.objetivos
    temario = data.temario
    encuadre = data.encuadre
    evaluaciones = data.evaluaciones

    beneficios = data.beneficios or ""
    expositiva = data.expositiva or {}
    demostrativa = data.demostrativa or {}
    dialogo = data.dialogo or {}
    cierre = data.cierre or {}

    def texto_seguro(valor):
        if valor is None:
            return ""
        if isinstance(valor, list):
            return "\n".join(str(v) for v in valor if v)
        return str(valor)

    def agregar_slide(titulo, contenido="", notas=""):
        slide_layout = prs.slide_layouts[1]  # título y contenido
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = texto_seguro(titulo)

        body = slide.placeholders[1]
        body.text = texto_seguro(contenido)

        if notas:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = texto_seguro(notas)

        return slide

    def agregar_portada():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = datos.nombreCurso or "Nombre del taller"

        subtitle = slide.placeholders[1]
        subtitle.text = (
            f"Instructor: {datos.instructor or ''}\n"
            f"Fecha: {datos.fecha or ''}\n"
            f"Lugar: {datos.lugar or ''}"
        )

    def lista_temario():
        lineas = []

        if getattr(temario, "u1", None):
            lineas.append("Unidad 1:")
            lineas.extend([f"• {t}" for t in temario.u1])

        if getattr(temario, "u2", None):
            lineas.append("\nUnidad 2:")
            lineas.extend([f"• {t}" for t in temario.u2])

        if getattr(temario, "u3", None):
            lineas.append("\nUnidad 3:")
            lineas.extend([f"• {t}" for t in temario.u3])

        return "\n".join(lineas)

    def texto_encuadre_reglas():
        reglas = []

        if hasattr(encuadre, "reglasTexto") and encuadre.reglasTexto:
            reglas.extend(encuadre.reglasTexto)

        if hasattr(encuadre, "otraRegla") and encuadre.otraRegla:
            reglas.append(encuadre.otraRegla)

        return "\n".join([f"• {r}" for r in reglas])

    def texto_encuadre_contrato():
        acuerdos = []

        if hasattr(encuadre, "acuerdosTexto") and encuadre.acuerdosTexto:
            acuerdos.extend(encuadre.acuerdosTexto)

        if hasattr(encuadre, "otroAcuerdo") and encuadre.otroAcuerdo:
            acuerdos.append(encuadre.otroAcuerdo)

        return "\n".join([f"• {a}" for a in acuerdos])

    def texto_preguntas_experiencia():
        preguntas = ""

        if hasattr(encuadre, "preguntas"):
            preguntas = encuadre.preguntas or ""

        return preguntas

    # 1. Portada
    agregar_portada()

    # 2. Presentación del ponente
    agregar_slide(
        "Presentación del Ponente",
        (
            f"Nombre del ponente: {datos.instructor or ''}\n\n"
            "Breve descripción de su formación profesional:\n\n"
            "Áreas de interés:\n\n"
            "Datos de contacto:\n\n"
            "Horario y lugar donde lo pueden localizar:\n\n"
            "Fotografía:"
        )
    )

    # 3. Lo que veremos en este curso
    agregar_slide(
        "Lo que veremos en este curso",
        f"Curso: {datos.nombreCurso or ''}\n\nPerfil de participantes:\n{datos.perfil or ''}"
    )

    # 4. Objetivo General
    agregar_slide(
        "Objetivo General",
        objetivos.general or ""
    )

    # 5. Objetivos Particulares
    agregar_slide(
        "Objetivos Particulares",
        (
            f"Cognitivo:\n{objetivos.cognitiva or ''}\n\n"
            f"Psicomotriz:\n{objetivos.psicomotriz or ''}\n\n"
            f"Afectivo:\n{objetivos.afectiva or ''}"
        )
    )

    # 6. Temario
    agregar_slide(
        "Temario",
        lista_temario()
    )

    # 7. Beneficios
    agregar_slide(
        "Beneficios de este curso",
        beneficios
    )

    # 8. Evaluación
    agregar_slide(
        "¿Cómo será la evaluación del curso?",
        (
            "Evaluación diagnóstica\n"
            f"Evaluación formativa: {evaluaciones.pctFormativa or 0}%\n"
            f"Evaluación sumativa: {evaluaciones.pctSumativa or 0}%\n"
            "Evaluación de reacción"
        )
    )

    # 9. Expectativas / preguntas
    agregar_slide(
        "¿Qué esperas de este curso?",
        texto_preguntas_experiencia()
    )

    # 10. Reglas
    agregar_slide(
        "Reglas del curso",
        texto_encuadre_reglas()
    )

    # 11. Contrato de aprendizaje
    agregar_slide(
        "Contrato de aprendizaje",
        texto_encuadre_contrato()
    )

    # 12. Evaluación diagnóstica
    agregar_slide(
        "Evaluación Diagnóstica",
        f"Instrucciones:\n{evaluaciones.instDiagnostica or ''}"
    )

    # 13. Desarrollo
    agregar_slide(
        "DESARROLLO DE LOS TEMAS",
        "A partir de aquí se deberán incluir todas las diapositivas de la técnica expositiva."
    )

    # 14. Fin técnica
    agregar_slide(
        "Fin de la técnica expositiva",
        ""
    )

    # 15. Actividad
    agregar_slide(
        "Actividad",
        (
            f"Nombre de la actividad:\n{demostrativa.get('actividad', '')}\n\n"
            f"Instrucciones:\n{demostrativa.get('ejemplos', '')}"
        )
    )

    # 16. Evaluación formativa
    agregar_slide(
        "Evaluación Formativa",
        (
            f"Instrucciones:\n{evaluaciones.instFormativa or ''}\n\n"
            f"Ponderación: {evaluaciones.pctFormativa or 0}%"
        )
    )

    # 17. Debate
    agregar_slide(
        "Debate",
        "Instrucciones:"
    )

    # 18. Diálogo / foro / mesa redonda
    agregar_slide(
        "Diálogo, foro, mesa redonda, etc.",
        f"Instrucciones:\n{dialogo.get('instrucciones', '')}"
    )

    # 19. Evaluación final
    agregar_slide(
        "Evaluación Final",
        (
            f"Instrucciones:\n{evaluaciones.instSumativa or ''}\n\n"
            f"Ponderación: {evaluaciones.pctSumativa or 0}%"
        )
    )

    # 20. Resumen
    agregar_slide(
        "Resumen",
        cierre.get('resumen', '')
    )

    # 21. Logro de expectativas
    agregar_slide(
        "Logro de expectativas",
        ""
    )

    # 22. Logro de objetivos
    agregar_slide(
        "Logro de objetivos",
        ""
    )

    # 23. Sugerencias
    agregar_slide(
        "Sugerencias de continuidad del aprendizaje",
        cierre.get('sugerencias', '')
    )

    # 24. Compromisos
    agregar_slide(
        "Compromisos de aplicación del aprendizaje",
        cierre.get('compromisos', '')
    )

    # 25. Encuesta
    agregar_slide(
        "Ayúdame a llenar la siguiente encuesta",
        f"Instrucciones:\n{evaluaciones.instReac or ''}"
    )

    # 26. Gracias
    agregar_slide(
        "¡GRACIAS!",
        ""
    )

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer.read()




# ─── Encabezado común para documentos ────────────────────────────────────────

def tabla_encabezado(doc: Document, datos: "DatosInfo") -> None:
    """Agrega la tabla de encabezado estándar EC0217 al documento."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    # 4 filas: se elimina la fila de "Perfil" y "Participantes"
    # Se agrega "Nombre del Participante" como línea en blanco para llenado manual
    tabla = doc.add_table(rows=4, cols=2)
    tabla.style = "Table Grid"

    def set_cell(row, col, label, value=""):
        c = tabla.rows[row].cells[col]
        c.text = ""
        p = c.paragraphs[0]
        run_label = p.add_run(label)
        run_label.bold = True
        if value:
            p.add_run(f" {value}")

    set_cell(0, 0, "Nombre del Curso/sesión:", datos.nombreCurso)
    set_cell(0, 1, "Lugar de Instrucción:", datos.lugar)
    set_cell(1, 0, "Nombre del facilitador/instructor/capacitador/formador:", datos.instructor)
    set_cell(1, 1, "Nombre del Diseñador:", datos.disenador)
    set_cell(2, 0, "Duración del curso:", f"{datos.duracion} minutos")
    set_cell(2, 1, "Fecha(s):", datos.fecha)
    # Fila 3: "Nombre del Participante" ocupa las dos columnas
    tabla.rows[3].cells[0].merge(tabla.rows[3].cells[1])
    set_cell(3, 0, "Nombre del Participante:", "")

    doc.add_paragraph()


def _docx_bytes(doc: Document) -> bytes:
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def _titulo(doc: Document, texto: str, subtitulo: str = "") -> None:
    h = doc.add_heading("", level=1)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = h.add_run(texto)
    r.font.color.rgb = RGBColor(0x1F, 0x3B, 0x6D)
    r.font.size = Pt(16)
    if subtitulo:
        h2 = doc.add_heading("", level=2)
        h2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = h2.add_run(subtitulo)
        r2.font.color.rgb = RGBColor(0x31, 0x4E, 0x7A)
        r2.font.size = Pt(13)


# ─── Documentos individuales ──────────────────────────────────────────────────

def generar_contrato_aprendizaje(data: "PlaneacionRequest") -> bytes:
    doc = Document()
    _titulo(doc, "CONTRATO DE APRENDIZAJE")
    tabla_encabezado(doc, data.datos)

    # ── Alcance e instrucciones del contrato ─────────────────────────────────
    doc.add_heading("Alcance e instrucciones del contrato de aprendizaje", level=2)
    doc.add_paragraph(
        "El presente contrato de aprendizaje establece los compromisos mutuos entre "
        "el instructor y los participantes para garantizar el logro de los objetivos "
        "del curso. Ambas partes se comprometen a cumplir con lo acordado durante el "
        "desarrollo de la sesión."
    )

    # Mostrar los acuerdos capturados en el wizard como alcance/instrucciones
    acuerdos = list(data.encuadre.acuerdosTexto or [])
    if data.encuadre.otroAcuerdo:
        acuerdos.append(data.encuadre.otroAcuerdo)
    if acuerdos:
        doc.add_paragraph()
        doc.add_paragraph("Los participantes acuerdan lo siguiente:")
        for i, acuerdo in enumerate(acuerdos, 1):
            doc.add_paragraph(f"{i}. {acuerdo}", style="List Number")

    doc.add_paragraph()

    # ── Compromisos del Instructor ────────────────────────────────────────────
    doc.add_heading("Compromisos del Instructor", level=2)
    compromisos_instructor = [
        "Me comprometo a que el curso sea dinámico y práctico.",
        "Me comprometo que los contenidos sean comprensibles y de fácil asimilación.",
        "Me comprometo a que se cumplan los objetivos del curso.",
        "Me comprometo a poder cumplir las expectativas del curso.",
        "Me comprometo a aclarar dudas respetando los tiempos del curso.",
        "Me comprometo a dejarte algo práctico para su vida profesional y laboral.",
    ]
    for i, c in enumerate(compromisos_instructor, 1):
        doc.add_paragraph(f"{i}. {c}", style="List Number")

    doc.add_paragraph()

    # ── Compromisos del Participante (líneas en blanco para llenar a mano) ───
    doc.add_heading("Compromisos del Participante", level=2)
    for i in range(1, 7):
        doc.add_paragraph(f"{i}. _______________________________________________", style="List Number")

    doc.add_paragraph()
    doc.add_paragraph("_________________________________          _______________________________")
    doc.add_paragraph("Nombre y Firma del Participante                        Firma del Instructor")

    return _docx_bytes(doc)


def generar_lista_requerimientos(data: "PlaneacionRequest") -> bytes:
    doc = Document()
    _titulo(doc, "LISTA DE VERIFICACIÓN DE REQUERIMIENTOS")
    tabla_encabezado(doc, data.datos)

    mat = data.materiales or {}

    # Mapeo: (título de sección, clave en materiales)
    secciones = [
        ("Instalaciones, mobiliario y su distribución",             mat.get("instalaciones", "")),
        ("Equipo de apoyo",                                         mat.get("equipo", "")),
        ("Materiales didácticos de apoyo y servicios",              mat.get("materialesDidacticos", "")),
        ("Requerimientos humanos",                                  mat.get("humanos", "")),
        ("Otros requerimientos",                                    mat.get("otros", "")),
        ("Medidas de salud/seguridad/higiene/protección civil",     mat.get("seguridad", "")),
    ]

    def _items(texto: str) -> list:
        """Convierte texto con viñetas en lista limpia de ítems."""
        if not texto:
            return []
        items = []
        for linea in texto.split("\n"):
            limpia = linea.strip().lstrip("•-–*").strip()
            if limpia:
                items.append(limpia)
        return items

    for titulo, contenido in secciones:
        doc.add_heading(titulo, level=2)

        items = _items(contenido)
        num_filas_datos = max(len(items), 3)   # mínimo 3 filas vacías
        t = doc.add_table(rows=num_filas_datos + 1, cols=4)
        t.style = "Table Grid"

        # Encabezados
        encabezados = ["No.", "Descripción", "Existe", "No existe"]
        for j, enc in enumerate(encabezados):
            t.rows[0].cells[j].paragraphs[0].add_run(enc).bold = True

        # Filas con datos
        for i, item in enumerate(items, start=1):
            t.rows[i].cells[0].text = str(i)
            t.rows[i].cells[1].text = item

        # Filas vacías restantes (solo número)
        for i in range(len(items) + 1, num_filas_datos + 1):
            t.rows[i].cells[0].text = str(i)

        doc.add_paragraph()

    return _docx_bytes(doc)


def generar_evaluacion_formativa_cotejo(data: "PlaneacionRequest") -> bytes:
    doc = Document()
    _titulo(doc, "EVALUACIÓN FORMATIVA", "LISTA DE COTEJO")
    tabla_encabezado(doc, data.datos)

    doc.add_heading("Instrucciones de aplicación", level=2)
    instrucciones = [
        ("Alcance/beneficio:", "La presente lista de cotejo tiene como propósito identificar la comprensión y avance logrado por el participante."),
        ("Indicaciones para el facilitador:", "Revise el documento entregado por el participante y marque con una X en SÍ cuando cumple con los reactivos solicitados, y en NO en caso contrario."),
        ("Condiciones de aplicación:", "Se aplicará en el desarrollo del curso durante la técnica demostrativa."),
        ("Tiempo para desarrollar la actividad:", "Considerar el mismo tiempo establecido en el documento de planeación."),
        ("Valor:", f"{data.evaluaciones.pctFormativa or data.evaluaciones.pctForm or 0}% de la calificación total."),
    ]
    for etiqueta, texto in instrucciones:
        p = doc.add_paragraph()
        p.add_run(etiqueta).bold = True
        p.add_run(f" {texto}")

    doc.add_paragraph()

    # ── Reactivos / criterios generados (tipoInstrumentoFormativa + instFormativa) ──
    reactivos_texto = (
        data.evaluaciones.instFormativa
        or data.evaluaciones.instForm
        or ""
    ).strip()
    if reactivos_texto:
        doc.add_heading("Criterios de evaluación", level=2)
        for linea in reactivos_texto.split("\n"):
            if linea.strip():
                doc.add_paragraph(linea.strip())
        doc.add_paragraph()

    t = doc.add_table(rows=7, cols=5)
    t.style = "Table Grid"
    encabezados = ["Reactivo", "Descripción", "SÍ", "NO", "Observaciones"]
    for j, enc in enumerate(encabezados):
        t.rows[0].cells[j].paragraphs[0].add_run(enc).bold = True
    for i in range(1, 7):
        t.rows[i].cells[0].text = str(i)

    return _docx_bytes(doc)


def generar_evaluacion_formativa_guia(data: "PlaneacionRequest") -> bytes:
    doc = Document()
    _titulo(doc, "EVALUACIÓN FORMATIVA", "GUÍA DE OBSERVACIÓN")
    tabla_encabezado(doc, data.datos)

    doc.add_heading("Instrucciones de aplicación", level=2)
    instrucciones = [
        ("Alcance/beneficio:", "La presente guía de observación tiene como propósito identificar la comprensión y correcto desempeño del participante."),
        ("Indicaciones para el facilitador:", "Observe cuidadosamente la ejecución de las actividades y marque con una X en SÍ cuando cumple con los reactivos solicitados, y en NO en caso contrario."),
        ("Condiciones de aplicación:", "Se aplicará en el desarrollo del curso durante la técnica demostrativa."),
        ("Tiempo para desarrollar la actividad:", "Considerar el mismo tiempo establecido en el documento de planeación."),
        ("Valor:", f"{data.evaluaciones.pctFormativa or data.evaluaciones.pctForm or 0}% de la calificación total."),
    ]
    for etiqueta, texto in instrucciones:
        p = doc.add_paragraph()
        p.add_run(etiqueta).bold = True
        p.add_run(f" {texto}")

    doc.add_paragraph()

    # ── Reactivos / criterios generados ──────────────────────────────────────
    reactivos_texto = (
        data.evaluaciones.instFormativa
        or data.evaluaciones.instForm
        or ""
    ).strip()
    if reactivos_texto:
        doc.add_heading("Criterios de evaluación", level=2)
        for linea in reactivos_texto.split("\n"):
            if linea.strip():
                doc.add_paragraph(linea.strip())
        doc.add_paragraph()

    t = doc.add_table(rows=7, cols=5)
    t.style = "Table Grid"
    encabezados = ["Reactivo", "Descripción", "SÍ", "NO", "Observaciones"]
    for j, enc in enumerate(encabezados):
        t.rows[0].cells[j].paragraphs[0].add_run(enc).bold = True
    for i in range(1, 7):
        t.rows[i].cells[0].text = str(i)

    return _docx_bytes(doc)


def generar_evaluacion_sumativa(data: "PlaneacionRequest") -> bytes:
    doc = Document()
    _titulo(doc, "EVALUACIÓN SUMATIVA", "CUESTIONARIO")
    tabla_encabezado(doc, data.datos)

    doc.add_heading("Instrucciones de aplicación", level=2)
    instrucciones = [
        ("Alcance/beneficio:", "El presente cuestionario tiene como propósito acreditar los aprendizajes adquiridos al finalizar el curso/taller."),
        ("Indicaciones para el facilitador:", f"La evaluación final se realiza de manera individual y tiene un valor del {data.evaluaciones.pctSumativa or data.evaluaciones.pctSuma or 0}%."),
        ("Condiciones de aplicación:", "Se aplicará al final del curso. El participante no deberá realizar consultas documentales o verbales."),
        ("Tiempo para desarrollar la actividad:", "Considerar el mismo tiempo establecido en el documento de planeación."),
        ("Instrumento:", data.evaluaciones.instSumativa or data.evaluaciones.instSuma or "Cuestionario"),
    ]
    for etiqueta, texto in instrucciones:
        p = doc.add_paragraph()
        p.add_run(etiqueta).bold = True
        p.add_run(f" {texto}")

    doc.add_paragraph()
    doc.add_paragraph("Nombre y firma del Participante: ________________________________________")
    doc.add_paragraph()

    # ── Reactivos generados por IA ────────────────────────────────────────────
    reactivos_texto = (
        data.evaluaciones.instSumativa
        or data.evaluaciones.instSuma
        or ""
    ).strip()

    if reactivos_texto:
        doc.add_heading("Cuestionario", level=2)
        for linea in reactivos_texto.split("\n"):
            if linea.strip():
                doc.add_paragraph(linea.strip())
    else:
        doc.add_heading("Cuestionario", level=2)
        for i in range(1, 11):
            doc.add_paragraph(f"{i}. __________________________________________________________")
            doc.add_paragraph()

    return _docx_bytes(doc)


def generar_evaluacion_reaccion(data: "PlaneacionRequest") -> bytes:
    doc = Document()
    _titulo(doc, "EVALUACIÓN DE REACCIÓN")
    tabla_encabezado(doc, data.datos)

    doc.add_paragraph("Instrucciones generales: Valore cada aspecto marcando con una X la columna que corresponda.")

    categorias = [
        ("De las características del evento", ["Se encuentra organizado", "Recibí atención cordial"]),
        ("Del contenido del curso", ["Se cubrió el objetivo general", "Se cubrieron los objetivos específicos", "Se realizaron actividades de aprendizaje"]),
        ("De las instalaciones", ["Están limpias", "Iluminación adecuada", "Cuenta con equipo adecuado", "Cuenta con mobiliario adecuado"]),
        ("Del desempeño del instructor", ["Organización de su trabajo", "Claridad al exponer", "Respetó a los participantes", "Promovió la participación"]),
        ("Del material didáctico", ["Estuvo organizado", "Presentaciones claras", "Estuvo entendible", "Está relacionado con el tema"]),
    ]
    escala = ["Excelente 10", "Bueno 8-9", "Regular 6-7", "Malo 5"]

    for categoria, items in categorias:
        doc.add_heading(categoria, level=2)
        t = doc.add_table(rows=len(items) + 1, cols=5)
        t.style = "Table Grid"
        encabezados = ["Aspecto"] + escala
        for j, enc in enumerate(encabezados):
            t.rows[0].cells[j].paragraphs[0].add_run(enc).bold = True
        for i, item in enumerate(items, 1):
            t.rows[i].cells[0].text = item
        doc.add_paragraph()

    doc.add_paragraph("En general, lo que más me gustó del taller fue: _________________________________")
    doc.add_paragraph("En general, lo que menos me gustó del taller fue: _______________________________")
    doc.add_paragraph("Comentarios y sugerencias: _________________________________________________")
    doc.add_paragraph()

    instReac_texto = (data.evaluaciones.instReac or "").strip()
    if instReac_texto:
        doc.add_heading("Preguntas adicionales", level=2)
        for linea in instReac_texto.split("\n"):
            if linea.strip():
                doc.add_paragraph(linea.strip())
        doc.add_paragraph()

    doc.add_paragraph("Gracias por tu participación")

    return _docx_bytes(doc)


def generar_lista_asistencia(data: "PlaneacionRequest") -> bytes:
    doc = Document()
    _titulo(doc, "LISTA DE ASISTENCIA")
    tabla_encabezado(doc, data.datos)

    t = doc.add_table(rows=21, cols=3)
    t.style = "Table Grid"
    encabezados = ["No.", "Nombre", "Firma"]
    for j, enc in enumerate(encabezados):
        t.rows[0].cells[j].paragraphs[0].add_run(enc).bold = True
    for i in range(1, 21):
        t.rows[i].cells[0].text = str(i)

    return _docx_bytes(doc)

# ─── Endpoints ────────────────────────────────────────────────────────────────

@app.get("/")
def home():
    return {"message": "SmartBuilder EC — API funcionando correctamente"}


@app.get("/health/integraciones")
async def health_integraciones(request: Request):
    """Verifica el estado de todas las integraciones externas. Solo superadmin."""
    import time as _time
    from datetime import datetime, timezone

    user = getattr(request.state, "user", None)
    if not user:
        raise HTTPException(status_code=401, detail="Autenticación requerida.")
    from database import get_supabase as _get_sb
    sb = _get_sb()
    prof = sb.table("profiles").select("rol").eq("id", user.get("sub")).single().execute()
    if not prof.data or prof.data.get("rol") != "super_admin":
        raise HTTPException(status_code=403, detail="Solo superadmin.")

    loop = asyncio.get_running_loop()

    async def _chk(key: str, fn):
        try:
            t0 = _time.time()
            extra = await loop.run_in_executor(None, fn)
            return key, {"status": "ok", "latency_ms": round((_time.time() - t0) * 1000), **(extra or {})}
        except Exception as e:
            return key, {"status": "error", "error": str(e)[:150]}

    def _test_claude():
        import anthropic as _ant
        key = os.getenv("ANTHROPIC_API_KEY", "")
        if not key:
            raise ValueError("ANTHROPIC_API_KEY no configurada")
        c = _ant.Anthropic(api_key=key)
        c.messages.create(model="claude-haiku-4-5-20251001", max_tokens=1,
                          messages=[{"role": "user", "content": "ping"}])
        return {"note": "Haiku + Sonnet disponibles"}

    def _test_openai():
        from openai import OpenAI as _OAI
        key = os.getenv("OPENAI_API_KEY", "")
        if not key:
            raise ValueError("OPENAI_API_KEY no configurada")
        c = _OAI(api_key=key)
        c.embeddings.create(model="text-embedding-3-small", input="ping")
        return {"model": "text-embedding-3-small"}

    def _test_supabase():
        r = sb.table("profiles").select("id", count="exact", head=True).execute()
        return {"perfiles": r.count}

    def _test_pgvector():
        zero = [0.0] * 1536
        sb.rpc("match_knowledge", {
            "query_embedding": zero,
            "similarity_threshold": 0.9999,
            "match_count": 1,
            "filtro_contexto": None,
        }).execute()
        r = sb.table("knowledge_faqs").select("id", count="exact", head=True).eq("activo", True).execute()
        return {"faqs_activas": r.count}

    def _test_resend():
        key = os.getenv("RESEND_API_KEY", "")
        if not key or not key.startswith("re_"):
            raise ValueError("RESEND_API_KEY no configurada")
        return {"from": os.getenv("RESEND_FROM_EMAIL", "configurado")}

    def _test_stripe():
        import stripe as _stripe
        key = os.getenv("STRIPE_SECRET_KEY", "")
        if not key or not key.startswith("sk_"):
            raise ValueError("STRIPE_SECRET_KEY no configurada")
        client = _stripe.StripeClient(key)
        client.balance.retrieve()
        return {"mode": "test" if "test" in key else "live"}

    def _test_docs():
        from docx import Document as _Doc
        from pptx import Presentation as _Prs
        import io
        buf = io.BytesIO()
        _Doc().save(buf)
        if buf.tell() == 0:
            raise ValueError("python-docx generó buffer vacío")
        buf2 = io.BytesIO()
        _Prs().save(buf2)
        if buf2.tell() == 0:
            raise ValueError("python-pptx generó buffer vacío")
        return {"docx": "ok", "pptx": "ok"}

    def _test_stripe_pagos():
        import stripe as _stripe
        key = os.getenv("STRIPE_SECRET_KEY", "")
        if not key or not key.startswith("sk_"):
            raise ValueError("STRIPE_SECRET_KEY no configurada")
        client = _stripe.StripeClient(key)
        sessions = client.checkout.sessions.list(params={"limit": 1, "status": "complete"})
        if not sessions.data:
            return {"ultimo": "sin pagos aún", "monto": None}
        s = sessions.data[0]
        diff = int(_time.time() - s.created)
        if diff < 3600:    hace = f"hace {diff // 60} min"
        elif diff < 86400: hace = f"hace {diff // 3600}h"
        else:              hace = f"hace {diff // 86400}d"
        monto = f"${(s.amount_total or 0) // 100:,.0f} MXN"
        return {"monto": monto, "hace": hace}

    def _test_tokens_ia():
        import httpx as _httpx
        from datetime import date as _date
        key = os.getenv("OPENAI_API_KEY", "")
        if not key:
            raise ValueError("OPENAI_API_KEY no configurada")
        today = _date.today().isoformat()
        r = _httpx.get(
            "https://api.openai.com/v1/usage",
            params={"date": today},
            headers={"Authorization": f"Bearer {key}"},
            timeout=8,
        )
        if r.status_code == 403:
            raise ValueError("Permiso 'Usage' no habilitado en esta API key — actívalo en dashboard.openai.com")
        if r.status_code != 200:
            raise ValueError(f"OpenAI Usage API respondió HTTP {r.status_code}")
        entries = r.json().get("data", [])
        tokens = sum(
            d.get("n_context_tokens_total", 0) + d.get("n_generated_tokens_total", 0)
            for d in entries
        )
        costo = round(tokens * 0.00000035, 4)
        return {"tokens": f"{tokens:,}", "costo": f"~${costo:.3f} USD"}

    checks = await asyncio.gather(
        _chk("claude",        _test_claude),
        _chk("openai",        _test_openai),
        _chk("supabase",      _test_supabase),
        _chk("pgvector",      _test_pgvector),
        _chk("resend",        _test_resend),
        _chk("stripe",        _test_stripe),
        _chk("docs",          _test_docs),
        _chk("stripe_pagos",  _test_stripe_pagos),
        _chk("tokens_ia",     _test_tokens_ia),
        _chk("tickets_kb",       lambda: _test_kb_pendientes("soporte_tickets",    "estado", "resuelto",  "ticket",     "sin resolver")),
        _chk("sugerencias_kb",   lambda: _test_kb_pendientes("soporte_sugerencias","estado", "pendiente", "sugerencia", "sin aplicar", invert=True)),
        _chk("vigencias_proximas", _test_vigencias_proximas),
        _chk("admins_sin_creditos", _test_admins_sin_creditos),
        _chk("usuarios_sin_plan",  _test_usuarios_sin_plan),
    )

def _test_vigencias_proximas():
    from database import get_supabase as _gsb
    from datetime import date, timedelta
    sb = _gsb()
    hoy    = date.today().isoformat()
    limite = (date.today() + timedelta(days=7)).isoformat()
    r = sb.table("profiles").select("id", count="exact", head=True)\
        .eq("rol", "admin").eq("activo", True)\
        .gte("vigencia_hasta", hoy).lte("vigencia_hasta", limite).execute()
    n = r.count or 0
    if n > 0:
        return {"status": "warning", "pendientes": n,
                "message": f"{n} admin{'s' if n != 1 else ''} con vigencia por vencer en ≤7 días"}
    return {"pendientes": 0}


def _test_admins_sin_creditos():
    from database import get_supabase as _gsb
    sb = _gsb()
    r = sb.table("profiles").select("id", count="exact", head=True)\
        .eq("rol", "admin").eq("activo", True).eq("credits", 0).execute()
    n = r.count or 0
    if n > 0:
        return {"status": "warning", "pendientes": n,
                "message": f"{n} admin{'s' if n != 1 else ''} con 0 créditos — no pueden registrar usuarios"}
    return {"pendientes": 0}


def _test_usuarios_sin_plan():
    from database import get_supabase as _gsb
    from datetime import datetime, timedelta, timezone
    sb = _gsb()
    hace_7 = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
    usuarios = sb.table("profiles").select("id")\
        .eq("rol", "user").eq("activo", True).lt("created_at", hace_7).limit(300).execute()
    ids = [u["id"] for u in (usuarios.data or [])]
    if not ids:
        return {"pendientes": 0}
    con_plan = sb.table("planeaciones").select("user_id").in_("user_id", ids).execute()
    con_ids  = {p["user_id"] for p in (con_plan.data or [])}
    n = len([i for i in ids if i not in con_ids])
    if n > 0:
        return {"status": "warning", "pendientes": n,
                "message": f"{n} usuario{'s' if n != 1 else ''} registrado{'s' if n != 1 else ''} hace +7 días sin ninguna planeación"}
    return {"pendientes": 0}


def _test_kb_pendientes(tabla, campo, valor, singular, descripcion, invert=False):
    from database import get_supabase as _gsb
    sb = _gsb()
    q = sb.table(tabla).select("id", count="exact", head=True)
    q = q.eq(campo, valor) if invert else q.neq(campo, valor)
    r = q.execute()
    n = r.count or 0
    if n > 0:
        msg = f"{n} {singular}{'s' if n != 1 else ''} {descripcion}"
        return {"status": "warning", "pendientes": n, "message": msg}
    return {"pendientes": 0}

    integraciones = dict(checks)
    todas_ok = all(v["status"] == "ok" for v in integraciones.values())

    return {
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "todas_ok": todas_ok,
        "integraciones": integraciones,
    }

@app.get("/perfil")
def get_perfil(request: Request):
    from database import get_supabase
    user_id = request.state.user.get("sub")
    sb = get_supabase()
    res = sb.table("profiles").select("id, nombre, apellido, rol, credits, activo").eq("id", user_id).single().execute()
    return res.data or {}


@app.post("/validate-token")
def validate_token(data: TokenRequest):
    SECRET_KEY = os.getenv("TOTP_SECRET", "JBSWY3DPEHPK3PXP")
    totp = pyotp.TOTP(SECRET_KEY, interval=30)
    if totp.verify(data.token):
        return {"status": "valid"}
    return {"status": "invalid"}


class CreateUserRequest(BaseModel):
    email: str
    nombre: str
    apellido: str
    admin_id: Optional[str] = None


@app.post("/admin/create-user", status_code=201)
def admin_create_user(data: CreateUserRequest, request: Request):
    from database import get_supabase

    caller_id = request.state.user.get("sub")
    sb = get_supabase()

    caller_res = sb.table("profiles").select("rol, id, nombre").eq("id", caller_id).single().execute()
    caller = caller_res.data or {}
    if caller.get("rol") not in ("admin", "super_admin"):
        raise HTTPException(status_code=403, detail="Solo admins pueden crear usuarios.")
    # Admin solo puede crear usuarios para sí mismo
    if caller.get("rol") == "admin" and caller.get("id") != data.admin_id:
        raise HTTPException(status_code=403, detail="No puedes crear usuarios para otro admin.")

    try:
        result = sb.auth.admin.create_user({
            "email": data.email,
            "email_confirm": True,
            "user_metadata": {
                "nombre": data.nombre,
                "apellido": data.apellido,
                "admin_id": data.admin_id,
            },
        })
        user_id = result.user.id

        # El trigger handle_new_user crea el perfil; aseguramos admin_id y rol
        sb.table("profiles").update({
            "nombre": data.nombre,
            "apellido": data.apellido,
            "admin_id": data.admin_id,
            "rol": "user",
            "activo": True,
        }).eq("id", user_id).execute()

        # Enviar email de configuración de contraseña
        frontend_url = os.getenv("FRONTEND_URL", "https://smartbuilderec.vercel.app")
        try:
            sb.auth.admin.generate_link({
                "type": "recovery",
                "email": data.email,
                "options": {"redirect_to": f"{frontend_url}/reset-password.html"},
            })
        except Exception:
            pass

        # Bienvenida al nuevo usuario
        try:
            from services.email_service import send_template as _send_tpl
            admin_nombre = caller.get("nombre") or "Tu administrador"
            if data.admin_id and caller.get("rol") == "super_admin":
                _ar = sb.table("profiles").select("nombre").eq("id", data.admin_id).single().execute()
                admin_nombre = (_ar.data or {}).get("nombre") or admin_nombre
            _send_tpl("bienvenida_user_codigo", data.email, {
                "nombre": data.nombre,
                "email": data.email,
                "nombre_admin": admin_nombre,
            })
        except Exception:
            pass

        # Alerta de créditos bajos cuando el admin llega a 1 crédito restante
        try:
            if data.admin_id:
                from services.email_service import send_template as _send_tpl
                _cr = sb.table("profiles").select("nombre, email, credits").eq("id", data.admin_id).single().execute()
                if _cr.data and _cr.data.get("credits") == 1 and _cr.data.get("email"):
                    _send_tpl("creditos_bajos_admin", _cr.data["email"], {
                        "nombre": _cr.data.get("nombre") or _cr.data["email"],
                        "email": _cr.data["email"],
                        "creditos_restantes": "1",
                    })
        except Exception:
            pass

        return {"id": user_id, "email": data.email}

    except Exception as e:
        err = str(e)
        # Registrar intento fallido en audit_logs para trazabilidad
        try:
            sb.table("audit_logs").insert({
                "actor_id":    caller_id,
                "actor_email": (sb.table("profiles").select("email").eq("id", caller_id).single().execute().data or {}).get("email", ""),
                "action":      "user_creation_failed",
                "target_email": data.email,
                "details":     {"error": err[:300]},
            }).execute()
        except Exception:
            pass
        if "already been registered" in err or "already exists" in err:
            raise HTTPException(status_code=409, detail="El correo ya está registrado.")
        raise HTTPException(status_code=500, detail=f"Error al crear usuario: {err}")


@app.delete("/admin/users/{user_id}", status_code=200)
def admin_delete_user(user_id: str, request: Request):
    from database import get_supabase

    caller_id = request.state.user.get("sub")
    sb = get_supabase()

    caller_res = sb.table("profiles").select("rol, id, email").eq("id", caller_id).single().execute()
    caller = caller_res.data or {}
    if caller.get("rol") not in ("admin", "super_admin"):
        raise HTTPException(status_code=403, detail="Sin permisos para eliminar usuarios.")

    # Obtener email del target ANTES de eliminar (para el audit log)
    target_res = sb.table("profiles").select("admin_id, rol, email").eq("id", user_id).single().execute()
    target = target_res.data or {}

    # Admin solo puede eliminar sus propios usuarios
    if caller.get("rol") == "admin":
        if target.get("admin_id") != caller_id:
            raise HTTPException(status_code=403, detail="No puedes eliminar usuarios de otro admin.")
        if target.get("rol") != "user":
            raise HTTPException(status_code=403, detail="Solo puedes eliminar usuarios regulares.")

    try:
        # Eliminar de auth.users → CASCADE a profiles (sin restaurar crédito al admin)
        sb.auth.admin.delete_user(user_id)

        # Audit log explícito: el trigger no puede obtener actor_id desde service_role
        try:
            sb.table("audit_logs").insert({
                "actor_id":    caller_id,
                "actor_email": caller.get("email", ""),
                "action":      "user_deleted",
                "target_id":   user_id,
                "target_email": target.get("email", ""),
                "details":     {}
            }).execute()
        except Exception:
            pass  # No bloquear el delete si el audit falla

        return {"deleted": user_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al eliminar usuario: {str(e)}")


@app.post("/evaluate")
def evaluate(data: EvaluationRequest):
    try:
        if data.tipo == "cognitiva":
            prompt = load_prompt(
                "evaluate_cognitiva_prompt.txt",
                texto=data.texto
            )

        elif data.tipo == "psicomotriz":
            prompt = load_prompt(
                "evaluate_psicomotriz_prompt.txt",
                texto=data.texto,
                objetivo_cognitivo=data.objetivo_cognitivo
            )

        elif data.tipo == "afectiva":
            prompt = load_prompt(
                "evaluate_afectiva_prompt.txt",
                texto=data.texto,
                objetivo_cognitivo=data.objetivo_cognitivo,
                objetivo_psicomotriz=data.objetivo_psicomotriz
            )

        else:
            raise HTTPException(
                status_code=400,
                detail="Tipo de objetivo no válido. Usa: cognitiva, psicomotriz o afectiva."
            )

        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        return json.loads(content.strip())

    except HTTPException:
        raise

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-general")
def generate_general(data: GeneralRequest):
    try:
        prompt = load_prompt(
            "general_prompt.txt",
            cognitiva=data.cognitiva,
            psicomotriz=data.psicomotriz,
            afectiva=data.afectiva,
        )
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0,
        )
        general = response.choices[0].message.content.strip()
        return {"general": general}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

@app.post("/generate-beneficios")
def generate_beneficios(data: BeneficiosRequest):
    try:
        prompt = load_prompt(
            "beneficios_prompt.txt",
            general=data.general,
            nombre=data.nombre
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
        )

        beneficios = response.choices[0].message.content.strip()

        return {"beneficios": beneficios}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



# ─── Generación del temario ───────────────────────────────────────

@app.post("/generate-temario")
def generate_temario(data: TemarioRequest):
    try:
        prompt = load_prompt(
            "temario_prompt.txt",
            nombre=data.nombre,
            general=data.general,
            cognitiva=data.cognitiva,
            psicomotriz=data.psicomotriz,
            afectiva=data.afectiva,
            beneficios=data.beneficios
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        resultado = json.loads(content.strip())

        return {
            "u1": resultado.get("u1", []),
            "u2": resultado.get("u2", []),
            "u3": resultado.get("u3", [])
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



@app.post("/generate-preguntas")
def generate_preguntas(data: PreguntasRequest):
    try:
        prompt = load_prompt(
            "preguntas_prompt.txt",
            nombre=data.nombre,
            perfil=data.perfil
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        resultado = json.loads(content.strip())

        return {
            "preguntas": resultado.get("preguntas", [])
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-doc/objetivos")
def generate_doc_objetivos(data: ObjetivosRequest):
    try:
        docx_bytes = crear_docx_objetivos(data)
        zip_bytes  = crear_zip_con_docx(docx_bytes)
        return StreamingResponse(
            io.BytesIO(zip_bytes),
            media_type="application/zip",
            headers={"Content-Disposition": "attachment; filename=objetivos_EC0217.zip"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


import re

def limpiar_nombre_archivo(texto: str) -> str:
    texto = str(texto or "").strip()
    texto = re.sub(r"\s+", "_", texto)
    texto = re.sub(r'[\\/:*?"<>|]', "", texto)
    texto = re.sub(r"_+", "_", texto)
    return texto or "Archivo"


@app.post("/generate-doc/planeacion")
def generate_doc_planeacion(data: PlaneacionRequest):
    try:
        payload = data.model_dump()
        nombre_curso = (data.datos.nombreCurso or "EC0217").replace(" ", "_")

        tipo_form = (data.evaluaciones.tipoInstrumentoFormativa or "").lower().strip()

        documentos = [
            ("01_Documento_de_Planeacion_EC0217.docx",  lambda: generar_planeacion_docx(payload)),
            ("02_Evaluacion_Diagnostica.docx",           lambda: generar_evaluacion_diagnostica(data)),
            ("05_Evaluacion_Sumativa.docx",              lambda: generar_evaluacion_sumativa(data)),
            ("06_Evaluacion_Reaccion.docx",              lambda: generar_evaluacion_reaccion(data)),
            ("07_Lista_de_Asistencia.docx",              lambda: generar_lista_asistencia(data)),
            ("08_Contrato_de_Aprendizaje.docx",          lambda: generar_contrato_aprendizaje(data)),
            ("09_Lista_de_Requerimientos.docx",          lambda: generar_lista_requerimientos(data)),
            ("10_Guia_de_Presentacion_EC0217.pptx",      lambda: generar_presentacion_curso(data)),
        ]

        if tipo_form == "guia_observacion":
            documentos.append(("03_Evaluacion_Formativa_Guia_Observacion.docx", lambda: generar_evaluacion_formativa_guia(data)))
        elif tipo_form == "lista_cotejo":
            documentos.append(("03_Evaluacion_Formativa_Cotejo.docx", lambda: generar_evaluacion_formativa_cotejo(data)))
        else:
            # Sin tipo definido: se incluyen ambos instrumentos
            documentos.append(("03_Evaluacion_Formativa_Cotejo.docx",           lambda: generar_evaluacion_formativa_cotejo(data)))
            documentos.append(("04_Evaluacion_Formativa_Guia_Observacion.docx", lambda: generar_evaluacion_formativa_guia(data)))


        # Guardar el payload como JSON para permitir reimportación en el wizard
        payload_exportable = data.model_dump()


        ev = payload_exportable.get("evaluaciones", {}) or {}

        payload_exportable["evaluaciones"] = {
            "pctDiagnostica": ev.get("pctDiagnostica", ev.get("pctDiag", 0)),
            "pctFormativa": ev.get("pctFormativa", ev.get("pctForm", 0)),
            "pctSumativa": ev.get("pctSumativa", ev.get("pctSuma", 0)),

            "instDiagnostica": ev.get("instDiagnostica", ev.get("instDiag", "")),
            "instFormativa": ev.get("instFormativa", ev.get("instForm", "")),
            "instSumativa": ev.get("instSumativa", ev.get("instSuma", "")),

            "instReac": ev.get("instReac", ""),
            "descripcionGeneral": ev.get("descripcionGeneral", ""),
            "tipoInstrumentoFormativa": ev.get("tipoInstrumentoFormativa", "")
        }
        


        payload_json_bytes = json.dumps(payload_exportable, ensure_ascii=False, indent=2).encode("utf-8")

        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
            # Incluir el JSON de datos para reimportación en el wizard
            instructor = payload_exportable.get("datos", {}).get("instructor", "Instructor")
            fecha = payload_exportable.get("datos", {}).get("fecha", "Sin_fecha")

            instructor = limpiar_nombre_archivo(instructor)
            fecha = limpiar_nombre_archivo(fecha)

            zf.writestr(f"{instructor}_{fecha}.json", payload_json_bytes)

            # Generar los 10 documentos en paralelo (reduce ~50 s → ~15 s)
            with ThreadPoolExecutor(max_workers=5) as executor:
                futures = {
                    executor.submit(generador): nombre_archivo
                    for nombre_archivo, generador in documentos
                }
                resultados = {}
                for future in as_completed(futures):
                    nombre_archivo = futures[future]
                    try:
                        resultados[nombre_archivo] = future.result()
                    except Exception as e_doc:
                        print(f"⚠️ Error generando {nombre_archivo}: {repr(e_doc)}")

            # Escribir al ZIP de forma secuencial (ZipFile no es thread-safe para escritura)
            for nombre_archivo, contenido in resultados.items():
                zf.writestr(nombre_archivo, contenido)

        zip_buffer.seek(0)

        return StreamingResponse(
            io.BytesIO(zip_buffer.read()),
            media_type="application/zip",
            headers={
                "Content-Disposition": f"attachment; filename=Expediente_{nombre_curso}.zip"
            }
        )

    except Exception as e:
        print("ERROR EN /generate-doc/planeacion:", repr(e))
        raise HTTPException(status_code=500, detail=str(e))
    
    

@app.post("/generate-expositiva")
def generate_expositiva(data: ExpositivaRequest):
    try:
        prompt = load_prompt(
            "expositiva_prompt.txt",
            campo=data.campo,
            nombreCurso=data.nombreCurso,
            perfil=data.perfil,
            objetivoCognitivo=data.objetivoCognitivo,
            objetivoGeneral=data.objetivoGeneral,
            temario=json.dumps(data.temario, ensure_ascii=False)
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        return json.loads(content.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    


@app.post("/generate-demostrativa")
def generate_demostrativa(data: DemostrativaRequest):
    try:
        prompt = load_prompt(
            "demostrativa_prompt.txt",
            campo=data.campo,
            nombreCurso=data.nombreCurso,
            perfil=data.perfil,
            objetivoPsicomotriz=data.objetivoPsicomotriz,
            objetivoGeneral=data.objetivoGeneral,
            temario=json.dumps(data.temario, ensure_ascii=False)
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        return json.loads(content.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    




@app.post("/generate-dialogo")
def generate_dialogo(data: DialogoRequest):
    try:
        prompt = load_prompt(
            "dialogo_prompt.txt",
            campo=data.campo,
            nombreCurso=data.nombreCurso,
            perfil=data.perfil,
            objetivoAfectivo=data.objetivoAfectivo,
            objetivoGeneral=data.objetivoGeneral,
            temario=json.dumps(data.temario, ensure_ascii=False)
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        return json.loads(content.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    


class ResumenRequest(BaseModel):
    nombreCurso:              str = ""
    objetivoGeneral:          str = ""
    objetivoCognitivo:        str = ""
    objetivoPsicomotriz:      str = ""
    objetivoAfectivo:         str = ""
    desarrolloExpositiva:     str = ""
    actividadDemostrativa:    str = ""
    instruccionesDialogo:     str = ""
    sugerenciasContinuidad:   str = ""
    referenciasBibliograficas: str = ""
    compromisos:              str = ""

@app.post("/generate-resumen")
def generate_resumen(data: ResumenRequest):
    try:
        prompt = load_prompt(
            "resumen_prompt.txt",
            nombreCurso              = data.nombreCurso,
            objetivoGeneral          = data.objetivoGeneral,
            objetivoCognitivo        = data.objetivoCognitivo,
            objetivoPsicomotriz      = data.objetivoPsicomotriz,
            objetivoAfectivo         = data.objetivoAfectivo,
            desarrolloExpositiva     = data.desarrolloExpositiva,
            actividadDemostrativa    = data.actividadDemostrativa,
            instruccionesDialogo     = data.instruccionesDialogo,
            sugerenciasContinuidad   = data.sugerenciasContinuidad,
            referenciasBibliograficas = data.referenciasBibliograficas,
            compromisos              = data.compromisos
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content_str = response.choices[0].message.content.strip()
        if content_str.startswith("```"):
            content_str = content_str.split("```")[1]
            if content_str.startswith("json"):
                content_str = content_str[4:]

        return json.loads(content_str.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


class CompromisosRequest(BaseModel):
    nombreCurso:        str = ""
    objetivoGeneral:    str = ""
    objetivoCognitivo:  str = ""
    objetivoPsicomotriz: str = ""
    objetivoAfectivo:   str = ""

@app.post("/generate-compromisos")
def generate_compromisos(data: CompromisosRequest):
    try:
        prompt = load_prompt(
            "compromisos_prompt.txt",
            nombreCurso         = data.nombreCurso,
            objetivoGeneral     = data.objetivoGeneral,
            objetivoCognitivo   = data.objetivoCognitivo,
            objetivoPsicomotriz = data.objetivoPsicomotriz,
            objetivoAfectivo    = data.objetivoAfectivo
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            response_format={"type": "json_object"}
        )

        content_str = response.choices[0].message.content.strip()
        if content_str.startswith("```"):
            content_str = content_str.split("```")[1]
            if content_str.startswith("json"):
                content_str = content_str[4:]

        return json.loads(content_str.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-cierre")
def generate_cierre(data: CierreRequest):
    try:
        prompt = load_prompt(
            "cierre_prompt.txt",
            nombreCurso=data.nombreCurso,
            objetivoGeneral=data.objetivoGeneral,
            objetivoCognitivo=data.objetivoCognitivo,
            objetivoPsicomotriz=data.objetivoPsicomotriz,
            objetivoAfectivo=data.objetivoAfectivo,
            desarrolloExpositiva=data.desarrolloExpositiva,
            actividadDemostrativa=data.actividadDemostrativa,
            instruccionesDialogo=data.instruccionesDialogo
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        return json.loads(content.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



@app.post("/generate-descripcion-general")
def generate_descripcion_general(data: DescripcionGeneralRequest):
    try:
        prompt = load_prompt(
            "descripcion_general_prompt.txt",
            cierre=data.cierre
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        return json.loads(content.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

    
@app.post("/generate-evaluacion-diagnostica")
def generate_evaluacion_diagnostica(data: EvaluacionIARequest):
    try:
        prompt = load_prompt(
            "evaluacion_diagnostica_prompt.txt",
            nombreCurso=data.nombreCurso,
            objetivoGeneral=data.objetivoGeneral,
            objetivoCognitivo=data.objetivoCognitivo,
            objetivoPsicomotriz=data.objetivoPsicomotriz,
            objetivoAfectivo=data.objetivoAfectivo
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        return json.loads(response.choices[0].message.content.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    



@app.post("/generate-evaluacion-sumativa")
def generate_evaluacion_sumativa(data: EvaluacionIARequest):
    try:
        prompt = load_prompt(
            "evaluacion_sumativa_prompt.txt",
            nombreCurso=data.nombreCurso,
            objetivoGeneral=data.objetivoGeneral,
            objetivoCognitivo=data.objetivoCognitivo,
            objetivoPsicomotriz=data.objetivoPsicomotriz,
            objetivoAfectivo=data.objetivoAfectivo
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        return json.loads(response.choices[0].message.content.strip())

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))





@app.post("/generate-materiales-clasificados")
async def generate_materiales_clasificados(payload: dict):
    integracion  = payload.get("integracion",  "")
    expositiva   = payload.get("expositiva",   "")
    demostrativa = payload.get("demostrativa", "")
    energizante  = payload.get("energizante",  "")
    dialogo      = payload.get("dialogo",      "")

    if not any([integracion, expositiva, demostrativa, energizante, dialogo]):
        raise HTTPException(status_code=400, detail="No hay materiales para clasificar.")

    try:
        prompt = load_prompt(
            "clasificar_materiales_prompt.txt",
            integracion=integracion   or "No especificado.",
            expositiva=expositiva     or "No especificado.",
            demostrativa=demostrativa or "No especificado.",
            energizante=energizante   or "No especificado.",
            dialogo=dialogo           or "No especificado.",
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()
        resultado = json.loads(content)

        campos = ["instalaciones", "equipo", "materialesDidacticos", "humanos", "otros", "seguridad"]
        return {c: resultado.get(c, "") for c in campos}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-formativa")
async def generate_formativa(payload: dict):
    nombre_curso = payload.get("nombreCurso", "")
    actividad = payload.get("actividad", "")

    if not actividad.strip():
        raise HTTPException(
            status_code=400,
            detail="Falta la actividad de la técnica demostrativa."
        )

    try:
        prompt = load_prompt(
            "evaluacion_formativa.txt",
            nombre_curso=nombre_curso,
            actividad=actividad
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[
                {
                    "role": "system",
                    "content": "Eres un experto en diseño de instrumentos de evaluación formativa por competencias."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.4,
            response_format={"type": "json_object"}
        )

        content = response.choices[0].message.content.strip()

        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]

        resultado = json.loads(content.strip())

        tipo = resultado.get("tipoInstrumento", "").strip()
        reactivos = resultado.get("reactivos", [])

        if not isinstance(reactivos, list):
            reactivos = []

        texto = "\n".join(
            f"{i + 1}. {reactivo}"
            for i, reactivo in enumerate(reactivos)
            if str(reactivo).strip()
        )

        return {
            "tipoInstrumento": tipo,
            "texto": texto
        }

    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error al generar evaluación formativa: {str(e)}"
        )
    


@app.post("/generate-materiales")
def generate_materiales(data: MaterialesRequest):
    try:
        prompts_materiales = {
            "integracion": "prompt_materiales_integracion.txt",
            "expositiva": "prompt_materiales_expositiva.txt",
            "demostrativa": "prompt_materiales_demostrativa.txt",
            "energizante": "prompt_materiales_energizante.txt",
            "dialogo": "prompt_materiales_dialogo.txt",
        }

        if data.tecnica not in prompts_materiales:
            raise HTTPException(
                status_code=400,
                detail="Técnica no válida."
            )

        archivo_prompt = prompts_materiales[data.tecnica]

        prompt = load_prompt(
            archivo_prompt,
            tecnica=data.tecnica,
            nombreCurso=data.nombreCurso,
            perfil=data.perfil,
            objetivoGeneral=data.objetivoGeneral,
            objetivos=json.dumps(data.objetivos, ensure_ascii=False, indent=2),
            temario=json.dumps(data.temario, ensure_ascii=False, indent=2),
            tecnicas=json.dumps(data.tecnicas, ensure_ascii=False, indent=2),
            expositiva=json.dumps(data.expositiva, ensure_ascii=False, indent=2),
            demostrativa=json.dumps(data.demostrativa, ensure_ascii=False, indent=2),
            dialogo=json.dumps(data.dialogo, ensure_ascii=False, indent=2),
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL_GRL,
            messages=[
                {
                    "role": "system",
                    "content": "Eres un experto en diseño instruccional bajo el estándar EC0217.01."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.4
        )

        texto = response.choices[0].message.content.strip()

        return {
            "texto": texto
        }

    except HTTPException:
        raise

    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error al generar materiales para {data.tecnica}: {str(e)}"
        )
    